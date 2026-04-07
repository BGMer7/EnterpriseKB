"""
文档解析器
支持PDF、Word、Excel、PPT、Markdown、图片等格式的解析
"""
from typing import List, Dict, Any, Optional
from io import BytesIO
from pathlib import Path
import logging
import os
from datetime import datetime

import pymupdf  # PyMuPDF
from docx import Document as DocxDocument
from openpyxl import load_workbook

logger = logging.getLogger(__name__)

# 尝试导入可选依赖
try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logger.warning("python-pptx not installed, PPT parsing disabled")

try:
    from pdf2image import convert_from_path
    from PIL import Image
    PDF2IMAGE_AVAILABLE = True
except ImportError:
    PDF2IMAGE_AVAILABLE = False
    logger.warning("pdf2image not installed, PDF image conversion disabled")

try:
    import pytesseract
    PYTESSERACT_AVAILABLE = True
except ImportError:
    PYTESSERACT_AVAILABLE = False
    logger.warning("pytesseract not installed, image OCR disabled")

try:
    from bs4 import BeautifulSoup
    BEAUTIFULSOUP_AVAILABLE = True
except ImportError:
    BEAUTIFULSOUP_AVAILABLE = False
    logger.warning("beautifulsoup4 not installed, HTML parsing disabled")

# PaddleOCR for PDF/image OCR
try:
    from paddleocr import PaddleOCR
    PADDLEOCR_AVAILABLE = True
except ImportError:
    PADDLEOCR_AVAILABLE = False
    logger.warning("paddleocr not installed, OCR disabled")


class DocumentParser:
    """
    文档解析器基类
    支持多种格式的文档解析
    """

    # 支持的格式映射
    SUPPORTED_TYPES = {
        # 原有格式
        "pdf": "_parse_pdf",
        "docx": "_parse_docx",
        "xlsx": "_parse_xlsx",
        "md": "_parse_markdown",
        "txt": "_parse_text",
        # 新增格式
        "pptx": "_parse_pptx",
        "ppt": "_parse_pptx",
        "doc": "_parse_doc",
        "html": "_parse_html",
        "htm": "_parse_html",
        "rtf": "_parse_rtf",
        # 图片格式
        "png": "_parse_image",
        "jpg": "_parse_image",
        "jpeg": "_parse_image",
        "gif": "_parse_image",
        "bmp": "_parse_image",
        "tiff": "_parse_image",
        "webp": "_parse_image",
    }

    @staticmethod
    def parse(file_path: str, file_type: str, use_ocr: bool = False, ocr_langs: str = 'ch,en') -> Dict[str, Any]:
        """
        解析文档

        Args:
            file_path: 文件路径
            file_type: 文件类型
            use_ocr: 是否使用OCR（对于PDF/图片有效）
            ocr_langs: OCR语言设置，默认为中英文 'ch,en'

        Returns:
            Dict: {
                "content": str,
                "pages": List[Dict],
                "metadata": Dict,
            }
        """
        file_type = file_type.lower()

        # 检查是否支持
        if file_type not in DocumentParser.SUPPORTED_TYPES:
            raise ValueError(f"Unsupported file type: {file_type}. Supported types: {', '.join(DocumentParser.SUPPORTED_TYPES.keys())}")

        # 获取解析方法
        method_name = DocumentParser.SUPPORTED_TYPES[file_type]
        parser_func = getattr(DocumentParser, method_name, None)

        if not parser_func:
            raise ValueError(f"Parser method not implemented for: {file_type}")

        # 对于PDF和图片，传递OCR参数
        if file_type in ["pdf", "png", "jpg", "jpeg", "gif", "bmp", "tiff", "webp"]:
            try:
                return parser_func(file_path, use_ocr=use_ocr, ocr_langs=ocr_langs)
            except TypeError:
                # 如果解析方法不支持OCR参数，回退到原方法
                return parser_func(file_path)

        # 调用解析方法
        try:
            return parser_func(file_path)
        except Exception as e:
            logger.error(f"Failed to parse {file_type} file {file_path}: {str(e)}")
            raise

    @staticmethod
    def parse_bytes(content: bytes, file_type: str) -> Dict[str, Any]:
        """
        解析字节内容

        Args:
            content: 文件字节内容
            file_type: 文件类型

        Returns:
            Dict: 解析结果
        """
        import tempfile
        import os

        suffix = f".{file_type}"
        with tempfile.NamedTemporaryFile(suffix=suffix, delete=False) as tmp:
            tmp.write(content)
            tmp_path = tmp.name

        try:
            result = DocumentParser.parse(tmp_path, file_type)
            return result
        finally:
            if os.path.exists(tmp_path):
                os.unlink(tmp_path)

    @staticmethod
    def _parse_pdf(file_path: str, use_ocr: bool = False, ocr_langs: str = 'ch,en') -> Dict[str, Any]:
        """
        解析PDF文档

        Args:
            file_path: PDF文件路径
            use_ocr: 是否使用OCR（用于扫描件）
            ocr_langs: OCR语言设置
        """
        doc = pymupdf.open(file_path)
        pages = []
        full_content = []
        ocr_used = False

        for page_num, page in enumerate(doc, 1):
            text = page.get_text()
            width = page.rect.width
            height = page.rect.height

            # 如果文本为空或过少，且启用OCR，则使用OCR
            if use_ocr and PADDLEOCR_AVAILABLE and (not text or len(text.strip()) < 50):
                logger.info(f"Page {page_num} has little text, using OCR...")
                try:
                    ocr_text = DocumentParser._ocr_page_with_paddle(file_path, page_num, ocr_langs)
                    if ocr_text:
                        text = ocr_text
                        ocr_used = True
                except Exception as e:
                    logger.warning(f"OCR failed for page {page_num}: {e}")

            # 提取图像和表格
            images = DocumentParser._extract_images_from_pdf_page(page, page_num)
            tables = DocumentParser._extract_tables_from_pdf_page(page, page_num)

            pages.append({
                "page_number": page_num,
                "content": text,
                "width": width,
                "height": height,
                "images": images,
                "tables": tables,
            })
            full_content.append(text)

        doc.close()

        # 获取文件信息
        file_stat = os.stat(file_path)

        metadata = {
            "page_count": len(pages),
            "format": "PDF",
            "file_name": os.path.basename(file_path),
            "file_size": file_stat.st_size,
            "file_created_date": datetime.fromtimestamp(file_stat.st_ctime).isoformat(),
            "file_modified_date": datetime.fromtimestamp(file_stat.st_mtime).isoformat(),
            "upload_date": datetime.now().isoformat(),
        }

        if use_ocr and ocr_used:
            metadata["ocr_used"] = True

        return {
            "content": "\n\n".join(full_content),
            "pages": pages,
            "metadata": metadata
        }

    @staticmethod
    def _extract_images_from_pdf_page(page, page_num: int) -> List[Dict[str, Any]]:
        """
        从PDF页面提取图像

        Args:
            page: PyMuPDF页面对象
            page_num: 页码

        Returns:
            List[Dict]: 图像信息列表
        """
        images = []
        image_list = page.get_images()

        for img_index, img in enumerate(image_list):
            try:
                xref = img[0]
                base_image = page.parent.extract_image(xref)

                image_info = {
                    "id": f"p{page_num}_img{img_index}",
                    "xref": xref,
                    "width": base_image.get("width"),
                    "height": base_image.get("height"),
                    "ext": base_image.get("ext"),
                    "size": len(base_image.get("image", b"")),
                    # 尝试获取图像位置
                    "position": DocumentParser._get_image_position(page, xref)
                }
                images.append(image_info)
            except Exception as e:
                logger.warning(f"Failed to extract image {img_index} from page {page_num}: {e}")

        return images

    @staticmethod
    def _get_image_position(page, xref: int) -> Optional[Dict[str, float]]:
        """获取图像在页面中的位置"""
        try:
            # 遍历页面元素查找图像
            for block in page.get_text("dict")["blocks"]:
                if block.get("type") == 1:  # 图像块
                    if block.get("xref") == xref:
                        bbox = block.get("bbox", [0, 0, 0, 0])
                        return {
                            "x": bbox[0],
                            "y": bbox[1],
                            "width": bbox[2] - bbox[0],
                            "height": bbox[3] - bbox[1]
                        }
        except Exception:
            pass
        return None

    @staticmethod
    def _extract_tables_from_pdf_page(page, page_num: int) -> List[Dict[str, Any]]:
        """
        从PDF页面提取表格

        Args:
            page: PyMuPDF页面对象
            page_num: 页码

        Returns:
            List[Dict]: 表格信息列表
        """
        tables = []

        try:
            # 使用表格提取功能（PyMuPDF 4.3.0+）
            tables_found = page.find_tables()

            if tables_found:
                for table_idx, table in enumerate(tables_found.tables):
                    try:
                        table_data = table.extract()

                        if table_data and len(table_data) > 0:
                            # 获取表格标题（查找表格上方的文本）
                            # 获取表格边界（bbox返回元组）
                            rect = table.bbox
                            if isinstance(rect, tuple):
                                x0, y0, x1, y1 = rect
                            else:
                                x0, y0, x1, y1 = rect.x0, rect.y0, rect.x1, rect.y1

                            # 查找表格上方的小文本作为标题
                            search_rect = pymupdf.Rect(x0, max(0, y0 - 50), x1, y0)

                            # 获取表格边界（bbox返回元组）
                            rect = table.bbox
                            if isinstance(rect, tuple):
                                x0, y0, x1, y1 = rect
                            else:
                                x0, y0, x1, y1 = rect.x0, rect.y0, rect.x1, rect.y1

                            # 简单实现：使用表格位置信息作为描述
                            table_info = {
                                "id": f"p{page_num}_table{table_idx}",
                                "row_count": len(table_data),
                                "col_count": len(table_data[0]) if table_data else 0,
                                "data": table_data,
                                "bbox": [x0, y0, x1, y1]
                            }
                            tables.append(table_info)
                    except Exception as e:
                        logger.warning(f"Failed to extract table {table_idx} from page {page_num}: {e}")
        except AttributeError:
            # find_tables方法不存在，使用备用方法
            logger.debug("find_tables not available, using fallback table extraction")
            tables = DocumentParser._extract_tables_fallback(page, page_num)
        except Exception as e:
            logger.warning(f"Table extraction failed for page {page_num}: {e}")

        return tables

    @staticmethod
    def _extract_tables_fallback(page, page_num: int) -> List[Dict[str, Any]]:
        """备用表格提取方法 - 基于文本对齐分析"""
        tables = []

        try:
            # 获取页面文本，按行分割
            text = page.get_text()
            lines = text.split("\n")

            # 简单的表格检测：检查是否有多个列对齐的行
            table_candidates = []
            current_table = []

            for line in lines:
                # 检测是否为表格行（多个制表符或空格分隔）
                if "\t" in line or "  " in line:
                    # 分割并清理
                    cells = [c.strip() for c in line.replace("\t", " | ").split(" | ") if c.strip()]
                    if len(cells) >= 2:
                        current_table.append(cells)
                else:
                    if len(current_table) >= 2:
                        table_candidates.append(current_table)
                    current_table = []

            if len(current_table) >= 2:
                table_candidates.append(current_table)

            # 创建表格对象
            for idx, table_data in enumerate(table_candidates):
                tables.append({
                    "id": f"p{page_num}_table{idx}",
                    "row_count": len(table_data),
                    "col_count": len(table_data[0]) if table_data else 0,
                    "data": table_data,
                    "extracted_by": "fallback"
                })

        except Exception as e:
            logger.warning(f"Fallback table extraction failed for page {page_num}: {e}")

        return tables

    @staticmethod
    def _ocr_page_with_paddle(file_path: str, page_num: int = 1, langs: str = 'ch,en') -> str:
        """
        使用PaddleOCR识别PDF指定页码的文字

        Args:
            file_path: PDF文件路径
            page_num: 页码（从1开始）
            langs: 语言设置

        Returns:
            识别出的文本
        """
        if not PADDLEOCR_AVAILABLE:
            raise ImportError("PaddleOCR is not installed")

        # 初始化OCR引擎（单次初始化，复用）
        if not hasattr(DocumentParser, '_ocr_engine'):
            logger.info("Initializing PaddleOCR engine...")
            DocumentParser._ocr_engine = PaddleOCR(
                use_angle_cls=True,
                lang='ch'
            )

        # 将PDF页面转换为图片
        doc = pymupdf.open(file_path)
        if page_num > len(doc):
            doc.close()
            return ""
        page = doc.load_page(page_num - 1)  # 页码从0开始
        pix = page.get_pixmap(matrix=pymupdf.Matrix(2, 2))  # 2x缩放提高清晰度

        # 保存为临时图片
        import tempfile
        import numpy as np

        img_data = pix.tobytes("png")
        img_array = np.frombuffer(img_data, dtype=np.uint8)
        img = img_array.reshape(pix.height, pix.width, 4)
        # 转为RGB（去除alpha通道）
        img = img[:, :, :3]

        doc.close()

        # OCR识别
        result = DocumentParser._ocr_engine.ocr(img, cls=True)

        if not result or not result[0]:
            return ""

        # 提取文本
        texts = []
        for line in result[0]:
            if line and len(line) >= 2:
                text = line[1][0]  # 识别出的文字
                texts.append(text)

        return "\n".join(texts)

    @staticmethod
    def _parse_docx(file_path: str) -> Dict[str, Any]:
        """
        解析Word文档(docx格式)
        """
        doc = DocxDocument(file_path)
        pages = []
        full_content = []

        # 收集所有段落和表格
        all_elements = []

        # 添加段落
        for para in doc.paragraphs:
            if para.text.strip():
                all_elements.append({
                    "type": "paragraph",
                    "text": para.text.strip(),
                    "style": str(para.style) if para.style else None
                })

        # 添加表格
        for table_idx, table in enumerate(doc.tables):
            table_data = []
            for row in table.rows:
                row_data = [cell.text.strip() for cell in row.cells]
                table_data.append(row_data)
            all_elements.append({
                "type": "table",
                "id": f"table_{table_idx}",
                "data": table_data
            })

        # 提取图像
        images = DocumentParser._extract_images_from_docx(doc)

        # 按页面组织（简单模拟，docx没有页概念）
        current_page = {
            "page_number": 1,
            "content": "",
            "elements": [],
            "images": images,
            "tables": []
        }

        sections = []
        table_count = 0
        for elem in all_elements:
            if elem["type"] == "paragraph":
                # 检测章节标题
                if elem.get("style") and "Heading" in str(elem["style"]):
                    sections.append({
                        "level": int(str(elem["style"]).split()[-1]) if elem["style"] else 1,
                        "text": elem["text"]
                    })
                current_page["content"] += elem["text"] + "\n"
                current_page["elements"].append(elem)
            elif elem["type"] == "table":
                # 表格内容转换为文本
                table_text = "\n".join([" | ".join(row) for row in elem["data"]])
                current_page["content"] += table_text + "\n"
                current_page["elements"].append(elem)
                # 添加表格信息
                current_page["tables"].append({
                    "id": elem.get("id", f"table_{table_count}"),
                    "row_count": len(elem["data"]),
                    "col_count": len(elem["data"][0]) if elem["data"] else 0,
                    "data": elem["data"]
                })
                table_count += 1

        current_page["sections"] = sections
        pages.append(current_page)
        full_content.append(current_page["content"])

        return {
            "content": "\n\n".join(full_content),
            "pages": pages,
            "metadata": {
                "page_count": 1,
                "format": "DOCX",
                "element_count": len(all_elements),
                "table_count": len(doc.tables),
                "image_count": len(images),
                "paragraph_count": len(doc.paragraphs)
            }
        }

    @staticmethod
    def _extract_images_from_docx(doc) -> List[Dict[str, Any]]:
        """
        从Word文档提取图像

        Args:
            doc: python-docx Document对象

        Returns:
            List[Dict]: 图像信息列表
        """
        images = []

        try:
            # 获取文档中的内联图像
            for rel in doc.part.rels.values():
                if "image" in rel.target_ref:
                    try:
                        image_part = rel.target_part
                        image_info = {
                            "id": rel.rId,
                            "content_type": image_part.content_type,
                            "size": len(image_part.blob),
                        }

                        # 尝试获取图像尺寸
                        if hasattr(image_part, "image"):
                            img = image_part.image
                            image_info["width"] = img.width if hasattr(img, "width") else None
                            image_info["height"] = img.height if hasattr(img, "height") else None

                        images.append(image_info)
                    except Exception as e:
                        logger.warning(f"Failed to extract image {rel.rId}: {e}")

        except Exception as e:
            logger.warning(f"Failed to extract images from docx: {e}")

        return images

    @staticmethod
    def _parse_doc(file_path: str) -> Dict[str, Any]:
        """
        解析老版本Word文档(.doc格式)
        尝试使用antiword或其他方式
        """
        import subprocess
        import os

        # 方法1: 尝试使用antiword
        try:
            result = subprocess.run(
                ["antiword", file_path],
                capture_output=True,
                text=True,
                timeout=30
            )
            if result.returncode == 0 and result.stdout:
                content = result.stdout
                return {
                    "content": content,
                    "pages": [{
                        "page_number": 1,
                        "content": content
                    }],
                    "metadata": {
                        "page_count": 1,
                        "format": "DOC",
                        "parser": "antiword"
                    }
                }
        except (FileNotFoundError, subprocess.TimeoutExpired):
            pass
        except Exception as e:
            logger.warning(f"antiword failed: {e}")

        # 方法2: 尝试使用pywin32 (Windows环境)
        try:
            import win32com.client
            word = win32com.client.Dispatch("Word.Application")
            word.Visible = False
            doc = word.Documents.Open(os.path.abspath(file_path))
            content = doc.Content.Text
            doc.Close(False)
            word.Quit()

            return {
                "content": content,
                "pages": [{
                    "page_number": 1,
                    "content": content
                }],
                "metadata": {
                    "page_count": 1,
                    "format": "DOC",
                    "parser": "win32com"
                }
            }
        except ImportError:
            pass
        except Exception as e:
            logger.warning(f"win32com failed: {e}")

        # 方法3: 使用python-docx2txt
        try:
            import docx2txt
            content = docx2txt.process(file_path)
            if content.strip():
                return {
                    "content": content,
                    "pages": [{
                        "page_number": 1,
                        "content": content
                    }],
                    "metadata": {
                        "page_count": 1,
                        "format": "DOC",
                        "parser": "docx2txt"
                    }
                }
        except ImportError:
            pass
        except Exception as e:
            logger.warning(f"docx2txt failed: {e}")

        raise ValueError(f"Unable to parse .doc file. Please install antiword, pywin32, or docx2txt: {file_path}")

    @staticmethod
    def _parse_pptx(file_path: str) -> Dict[str, Any]:
        """
        解析PPT/PPTX文档
        """
        if not PPTX_AVAILABLE:
            raise ImportError("python-pptx is not installed. Install with: pip install python-pptx")

        prs = Presentation(file_path)
        pages = []
        full_content = []

        # 幻灯片元数据
        slide_titles = []

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_content = ""
            elements = []

            # 提取标题
            title = None
            if slide.shapes.title:
                title_text = slide.shapes.title.text.strip()
                if title_text:
                    title = title_text
                    slide_content += f"标题: {title_text}\n"
                    elements.append({"type": "title", "text": title_text})

            slide_titles.append(title or f"Slide {slide_num}")

            # 提取图像
            slide_images = []
            for shape_idx, shape in enumerate(slide.shapes):
                # 提取图片
                if hasattr(shape, "image"):
                    try:
                        img = shape.image
                        image_info = {
                            "id": f"slide{slide_num}_img{shape_idx}",
                            "content_type": img.content_type if hasattr(img, "content_type") else None,
                            "width": img.size[0] if hasattr(img, "size") else None,
                            "height": img.size[1] if hasattr(img, "size") else None,
                        }
                        slide_images.append(image_info)
                    except Exception as e:
                        logger.warning(f"Failed to extract image from slide {slide_num}: {e}")

                if hasattr(shape, "text") and shape.text.strip():
                    if shape != slide.shapes.title:  # 跳过标题（已处理）
                        slide_content += shape.text.strip() + "\n"
                        elements.append({"type": "text", "text": shape.text.strip()})

            # 提取所有文本框内容
            slide_tables = []
            for shape in slide.shapes:
                # 提取表格内容
                if shape.shape_type == 19:  # MSO_SHAPE_TYPE.TABLE
                    table_data = []
                    for row in shape.table.rows:
                        row_data = [cell.text.strip() for cell in row.cells]
                        table_data.append(row_data)
                    table_text = "\n".join([" | ".join(row) for row in table_data])
                    slide_content += f"[表格]\n{table_text}\n"
                    elements.append({"type": "table", "data": table_data})
                    slide_tables.append({
                        "id": f"slide{slide_num}_table{len(slide_tables)}",
                        "row_count": len(table_data),
                        "col_count": len(table_data[0]) if table_data else 0,
                        "data": table_data
                    })

            pages.append({
                "page_number": slide_num,
                "title": title,
                "content": slide_content.strip(),
                "elements": elements,
                "images": slide_images,
                "tables": slide_tables
            })
            full_content.append(slide_content.strip())

        return {
            "content": "\n\n---\n\n".join(full_content),
            "pages": pages,
            "metadata": {
                "page_count": len(pages),
                "format": "PPTX",
                "slide_titles": slide_titles,
                "title": slide_titles[0] if slide_titles else None
            }
        }

    @staticmethod
    def _parse_xlsx(file_path: str) -> Dict[str, Any]:
        """
        解析Excel文档
        """
        wb = load_workbook(file_path, data_only=True)
        pages = []
        full_content = []

        for sheet_num, sheet_name in enumerate(wb.sheetnames):
            sheet = wb[sheet_name]
            sheet_content = []

            for row in sheet.iter_rows(values_only=True):
                row_text = "\t".join([str(cell) if cell is not None else "" for cell in row])
                sheet_content.append(row_text)

            pages.append({
                "page_number": sheet_num + 1,
                "sheet_name": sheet_name,
                "content": "\n".join(sheet_content)
            })

            full_content.extend(sheet_content)

        return {
            "content": "\n\n".join(full_content),
            "pages": pages,
            "metadata": {
                "page_count": len(pages),
                "format": "XLSX",
                "sheet_count": len(wb.sheetnames),
                "sheet_names": wb.sheetnames
            }
        }

    @staticmethod
    def _parse_markdown(file_path: str) -> Dict[str, Any]:
        """
        解析Markdown文档
        """
        with open(file_path, "r", encoding="utf-8") as f:
            md_content = f.read()

        pages = [{
            "page_number": 1,
            "content": md_content
        }]

        return {
            "content": md_content,
            "pages": pages,
            "metadata": {
                "page_count": 1,
                "format": "Markdown",
                "file_type": "md"
            }
        }

    @staticmethod
    def _parse_html(file_path: str) -> Dict[str, Any]:
        """
        解析HTML文档
        """
        if not BEAUTIFULSOUP_AVAILABLE:
            raise ImportError("beautifulsoup4 is not installed. Install with: pip install beautifulsoup4")

        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            html_content = f.read()

        # 解析HTML
        soup = BeautifulSoup(html_content, 'html.parser')

        # 移除script和style标签
        for tag in soup(["script", "style"]):
            tag.decompose()

        # 获取文本
        text_content = soup.get_text()

        # 清理空白
        lines = (line.strip() for line in text_content.splitlines())
        chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
        text_content = ' '.join(chunk for chunk in chunks if chunk)

        # 提取标题
        title = soup.title.string if soup.title else None

        # 提取链接
        links = []
        for a in soup.find_all('a', href=True):
            links.append({
                "text": a.get_text(strip=True),
                "href": a['href']
            })

        pages = [{
            "page_number": 1,
            "content": text_content,
            "title": title,
            "links": links[:20]  # 限制数量
        }]

        return {
            "content": text_content,
            "pages": pages,
            "metadata": {
                "page_count": 1,
                "format": "HTML",
                "title": title,
                "link_count": len(links)
            }
        }

    @staticmethod
    def _parse_rtf(file_path: str) -> Dict[str, Any]:
        """
        解析RTF文档
        """
        import re

        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            rtf_content = f.read()

        # 简单RTF解析：去除控制符，提取纯文本
        # 移除RTF控制词（如 \par, \n, \tab 等）
        text = re.sub(r'\\[a-z]+\d*\s?', ' ', rtf_content)
        # 移除大括号
        text = re.sub(r'[\{\}]', '', text)
        # 移除特殊字符
        text = re.sub(r'[\x00-\x1f\x7f-\x9f]', '', text)
        # 清理空白
        text = ' '.join(text.split())

        pages = [{
            "page_number": 1,
            "content": text
        }]

        return {
            "content": text,
            "pages": pages,
            "metadata": {
                "page_count": 1,
                "format": "RTF"
            }
        }

    @staticmethod
    def _parse_image(file_path: str, use_ocr: bool = False, ocr_langs: str = 'ch,en') -> Dict[str, Any]:
        """
        解析图片文档（使用OCR提取文字）

        Args:
            file_path: 图片文件路径
            use_ocr: 是否使用OCR
            ocr_langs: OCR语言设置
        """
        from PIL import Image
        import numpy as np

        # 打开图片
        image = Image.open(file_path)

        # 获取图片信息
        width, height = image.size
        format_type = image.format

        text = ""
        ocr_method = None

        # 优先使用PaddleOCR
        if PADDLEOCR_AVAILABLE:
            try:
                logger.info("Using PaddleOCR for image...")
                # 转换为numpy数组
                img_array = np.array(image)
                if img_array.ndim == 2:
                    # 灰度图
                    img = img_array
                elif img_array.ndim == 3:
                    if img_array.shape[2] == 4:
                        # RGBA -> RGB
                        img = img_array[:, :, :3]
                    else:
                        img = img_array
                else:
                    img = img_array

                # 初始化OCR引擎
                if not hasattr(DocumentParser, '_ocr_engine'):
                    DocumentParser._ocr_engine = PaddleOCR(
                        use_angle_cls=True,
                        lang='ch'
                    )

                result = DocumentParser._ocr_engine.ocr(img, cls=True)
                if result and result[0]:
                    texts = []
                    for line in result[0]:
                        if line and len(line) >= 2:
                            texts.append(line[1][0])
                    text = "\n".join(texts)
                ocr_method = "paddleocr"
            except Exception as e:
                logger.warning(f"PaddleOCR failed: {e}")

        # 回退到pytesseract
        if not text and PYTESSERACT_AVAILABLE:
            try:
                logger.info("Using pytesseract for image...")
                # 检查Tesseract是否可用
                pytesseract.get_tesseract_version()
                text = pytesseract.image_to_string(image, lang='chi_sim+eng')
                ocr_method = "pytesseract"
            except Exception as e:
                logger.warning(f"pytesseract failed: {e}")

        if not text:
            raise RuntimeError("No OCR method available. Please install PaddleOCR or pytesseract with Tesseract")

        # 清理文本
        text = text.strip()

        pages = [{
            "page_number": 1,
            "content": text,
            "width": width,
            "height": height
        }]

        return {
            "content": text,
            "pages": pages,
            "metadata": {
                "page_count": 1,
                "format": format_type or "IMAGE",
                "width": width,
                "height": height,
                "has_text": bool(text.strip()),
                "ocr_method": ocr_method,
                "char_count": len(text)
            }
        }

    @staticmethod
    def _parse_text(file_path: str) -> Dict[str, Any]:
        """
        解析纯文本文档
        """
        # 尝试多种编码
        encodings = ['utf-8', 'gbk', 'gb2312', 'latin-1']

        text_content = None
        used_encoding = None

        for encoding in encodings:
            try:
                with open(file_path, "r", encoding=encoding) as f:
                    text_content = f.read()
                used_encoding = encoding
                break
            except UnicodeDecodeError:
                continue

        if text_content is None:
            # 最后尝试二进制读取
            with open(file_path, "rb") as f:
                text_content = f.read().decode('utf-8', errors='ignore')
            used_encoding = 'utf-8 (with errors)'

        pages = [{
            "page_number": 1,
            "content": text_content
        }]

        return {
            "content": text_content,
            "pages": pages,
            "metadata": {
                "page_count": 1,
                "format": "TXT",
                "encoding": used_encoding,
                "char_count": len(text_content)
            }
        }


def parse_document(file_path: str, file_type: str) -> Dict[str, Any]:
    """
    解析文档（便捷函数）

    Args:
        file_path: 文件路径
        file_type: 文件类型

    Returns:
        Dict: 解析结果
    """
    return DocumentParser.parse(file_path, file_type)


def parse_document_from_bytes(content: bytes, file_type: str) -> Dict[str, Any]:
    """
    解析文档字节内容（便捷函数）

    Args:
        content: 文件字节内容
        file_type: 文件类型

    Returns:
        Dict: 解析结果
    """
    return DocumentParser.parse_bytes(content, file_type)
