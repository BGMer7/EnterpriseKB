"""
文档解析演示脚本
用于测试各种文档格式的解析效果
"""
import os
import sys
import tempfile
from pathlib import Path

# 添加项目根目录到路径
sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from app.processors.parser import DocumentParser


def create_sample_files(temp_dir: str) -> dict:
    """创建示例文件"""
    files = {}

    # 1. 文本文件
    txt_path = os.path.join(temp_dir, "sample.txt")
    with open(txt_path, "w", encoding="utf-8") as f:
        f.write("""这是测试文本文件
包含多行内容
第三行内容
""")
    files["txt"] = txt_path

    # 2. Markdown文件
    md_path = os.path.join(temp_dir, "sample.md")
    with open(md_path, "w", encoding="utf-8") as f:
        f.write("""# 标题一

这是正文内容。

## 二级标题

- 列表项1
- 列表项2
- 列表项3

### 代码示例

```python
def hello():
    print("Hello World")
```
""")
    files["md"] = md_path

    # 3. HTML文件
    html_path = os.path.join(temp_dir, "sample.html")
    with open(html_path, "w", encoding="utf-8") as f:
        f.write("""<!DOCTYPE html>
<html>
<head>
    <title>测试页面</title>
</head>
<body>
    <h1>欢迎访问</h1>
    <p>这是一个测试HTML页面</p>
    <a href="http://example.com">链接</a>
    <ul>
        <li>项目1</li>
        <li>项目2</li>
    </ul>
</body>
</html>
""")
    files["html"] = html_path

    # 4. RTF文件
    rtf_path = os.path.join(temp_dir, "sample.rtf")
    with open(rtf_path, "w", encoding="utf-8") as f:
        f.write(r"""{\rtf1\ansi\deff0
{\fonttbl{\f0 Courier;}}
\pard\f0\fs20 This is {\b bold} text.\par
This is normal text.\par
}
""")
    files["rtf"] = rtf_path

    # 5. Word文件 (docx)
    try:
        from docx import Document
        docx_path = os.path.join(temp_dir, "sample.docx")
        doc = Document()
        doc.add_heading("Word文档测试", 0)
        doc.add_paragraph("这是第一段内容，包含一些文字。")
        doc.add_paragraph("这是第二段内容。")

        # 添加表格
        table = doc.add_table(rows=3, cols=2)
        table.cell(0, 0).text = "姓名"
        table.cell(0, 1).text = "年龄"
        table.cell(1, 0).text = "张三"
        table.cell(1, 1).text = "25"
        table.cell(2, 0).text = "李四"
        table.cell(2, 1).text = "30"

        doc.add_heading("二级标题", 1)
        doc.add_paragraph("更多内容...")
        doc.save(docx_path)
        files["docx"] = docx_path
    except ImportError:
        print("⚠️ python-docx 未安装，跳过docx")

    # 6. Excel文件 (xlsx)
    try:
        from openpyxl import Workbook
        xlsx_path = os.path.join(temp_dir, "sample.xlsx")
        wb = Workbook()

        ws1 = wb.active
        ws1.title = "员工表"
        ws1["A1"] = "姓名"
        ws1["B1"] = "部门"
        ws1["C1"] = "薪资"
        ws1["A2"] = "张三"
        ws1["B2"] = "技术部"
        ws1["C2"] = 10000
        ws1["A3"] = "李四"
        ws1["B3"] = "人事部"
        ws1["C3"] = 8000

        ws2 = wb.create_sheet("部门表")
        ws2["A1"] = "部门"
        ws2["B1"] = "负责人"
        ws2["A2"] = "技术部"
        ws2["B2"] = "王五"

        wb.save(xlsx_path)
        files["xlsx"] = xlsx_path
    except ImportError:
        print("⚠️ openpyxl 未安装，跳过xlsx")

    # 7. PPT文件 (pptx)
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt

        pptx_path = os.path.join(temp_dir, "sample.pptx")
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)

        # 标题页
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = "演示文稿标题"
        subtitle.text = "副标题：测试内容"

        # 内容页
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        title_shape.text = "要点列表"
        tf = body_shape.text_frame
        tf.text = "第一点"
        p = tf.add_paragraph()
        p.text = "第二点"
        p.level = 1
        p = tf.add_paragraph()
        p.text = "第三点"
        p.level = 1

        # 表格页 - 使用标题和内容布局
        title_only_layout = prs.slide_layouts[5]  # Title Only
        slide = prs.slides.add_slide(title_only_layout)
        shapes = slide.shapes
        if shapes.title:
            shapes.title.text = "数据表格"

        rows, cols = 3, 3
        left = Inches(2)
        top = Inches(2)
        width = Inches(6)
        height = Inches(1.5)
        table = shapes.add_table(rows, cols, left, top, width, height).table

        table.cell(0, 0).text = "项目"
        table.cell(0, 1).text = "Q1"
        table.cell(0, 2).text = "Q2"
        table.cell(1, 0).text = "销量"
        table.cell(1, 1).text = "100"
        table.cell(1, 2).text = "150"
        table.cell(2, 0).text = "营收"
        table.cell(2, 1).text = "10000"
        table.cell(2, 2).text = "15000"

        prs.save(pptx_path)
        files["pptx"] = pptx_path
    except ImportError:
        print("⚠️ python-pptx 未安装，跳过pptx")

    # 8. PDF文件 (创建简单的文本PDF)
    try:
        import pymupdf
        pdf_path = os.path.join(temp_dir, "sample.pdf")
        doc = pymupdf.open()
        page = doc.new_page(width=595, height=842)  # A4大小

        # 写入文本
        text = """PDF测试文档

这是第一页的内容。

标题：测试PDF解析

正文内容：
1. 第一点
2. 第二点
3. 第三点
"""
        page.insert_text((50, 50), text, fontsize=12)

        # 添加第二页
        page2 = doc.new_page(width=595, height=842)
        page2.insert_text((50, 50), "这是第二页的内容\n\nPDF支持多页解析", fontsize=12)

        doc.save(pdf_path)
        doc.close()
        files["pdf"] = pdf_path
    except ImportError:
        print("⚠️ PyMuPDF 未安装，跳过pdf")

    return files


def print_result(file_type: str, result: dict):
    """打印解析结果"""
    print(f"\n{'='*60}")
    print(f"📄 格式: {file_type.upper()}")
    print(f"{'='*60}")
    print(f"📊 元数据: {result.get('metadata', {})}")
    print(f"📝 内容预览 (前500字符):")
    print("-" * 40)
    content = result.get('content', '')
    print(content[:500] if len(content) > 500 else content)
    if len(content) > 500:
        print(f"... (共 {len(content)} 字符)")
    print("-" * 40)


def main():
    """主函数"""
    print("🚀 文档解析演示程序")
    print("=" * 60)

    # 创建临时目录
    with tempfile.TemporaryDirectory() as temp_dir:
        print(f"📁 临时目录: {temp_dir}")
        print("\n📋 正在创建示例文件...")

        # 创建示例文件
        files = create_sample_files(temp_dir)

        print(f"✅ 创建了 {len(files)} 个示例文件: {list(files.keys())}")

        # 解析每个文件
        print("\n" + "="*60)
        print("🔍 开始解析...")
        print("="*60)

        for file_type, file_path in files.items():
            if not os.path.exists(file_path):
                continue

            try:
                print(f"\n⏳ 正在解析 {file_type} 文件...")
                result = DocumentParser.parse(file_path, file_type)
                print_result(file_type, result)
                print(f"✅ {file_type} 解析成功!")
            except Exception as e:
                print(f"❌ {file_type} 解析失败: {str(e)}")

        print("\n" + "="*60)
        print("🎉 演示完成!")
        print("="*60)


if __name__ == "__main__":
    main()
